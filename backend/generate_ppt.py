import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def apply_dark_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(14, 18, 20)  # Very dark, almost black bg

def format_title(title_shape, text):
    title_shape.text = text
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Georgia" 
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(78, 205, 196) # EmpathEase Teal

def format_body(body_shape, points):
    tf = body_shape.text_frame
    tf.clear()
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Pt(18)
        
        # Check if the line is a major point or sub-point
        # Major points get bigger text, sub-points get smaller and are indented
        if not point.startswith("  "):
            p.level = 0
            for run in p.runs:
                run.font.name = "Arial"
                run.font.size = Pt(22)
                run.font.color.rgb = RGBColor(238, 242, 240)
                # Highlighting specific impactful keywords
                if any(keyword in point for keyword in ["93%", "Zero", "Hard-Coded", "NPU", "RAG", "Real-Time"]):
                    run.font.bold = True 
        else:
            p.level = 1
            for run in p.runs:
                run.font.name = "Arial"
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(180, 190, 190)

def create_presentation():
    prs = Presentation()
    
    # SLIDE 1: Title Slide (High Impact)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_dark_theme(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "EmpathEase v1.1"
    subtitle.text = "The Edge of Empathy.\nMultimodal AI Therapy on AMD Architecture."
    
    title.text_frame.paragraphs[0].runs[0].font.name = "Georgia"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(64)
    title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(78, 205, 196)
    
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(26)
            run.font.color.rgb = RGBColor(200, 210, 210)

    # Content for remaining slides - Rewritten for maximum impact
    slides_data = [
        {
            "title": "The Status Quo is Failing Us",
            "content": [
                "Mental healthcare is broken, inaccessible, and reactive.",
                "  • India has only 0.75 psychiatrists per 100,000 people.",
                "  • Seeking help is stigmatized; people wait until absolute crisis.",
                "Current 'AI Therapists' are just glorified chatbots.",
                "  • They only read text.",
                "  • They miss 93% of human communication: micro-expressions and vocal tone."
            ]
        },
        {
            "title": "Enter EmpathEase: Multimodal Empathy",
            "content": [
                "We don't just read messages. We see, hear, and understand.",
                "  • Real-Time Fusion: Synchronous analysis of facial expressions, vocal tone, and text semantics.",
                "  • Culturally Native: Fluent in Hinglish, grounded in Indian socio-clinical contexts.",
                "  • The 'Antidote' UI: The interface breathes and changes color to actively down-regulate distress (Chromatherapy)."
            ]
        },
        {
            "title": "The Tech: How We Build Empathy",
            "content": [
                "A continuous, low-latency WebSocket stream powers our pipeline:",
                "1. Edge Capture: React/Vite client running 60fps frame capture & audio chunking.",
                "2. The Fusion Layer: Reconciles conflicting signals (e.g., smiling while saying 'I want to give up').",
                "3. Therapy Engine: Ingests the 'Fused Emotional State' to dynamically alter LLM system prompts mid-session."
            ]
        },
        {
            "title": "RAG: Beyond Generic Chat",
            "content": [
                "EmpathEase refuses to give dangerous, generic advice.",
                "  • Clinical Guardrails: Vector DB loaded with safe psychoeducation pathways.",
                "  • Cultural Nuance: Automatically pulls relevant idioms and culturally-safe responses.",
                "  • Dynamic Grounding: Stops LLM hallucinations in sensitive mental health scenarios."
            ]
        },
        {
            "title": "Zero-Trust Privacy & Safety",
            "content": [
                "In mental health, privacy isn't a feature; it's the foundation.",
                "  • Ephemeral Sessions: Zero server-side storage of PII, audio, or video.",
                "  • The Crisis Override: If severe distress is detected, the LLM is hard-killed.",
                "  • Human Handoff: Immediate, un-bypassable injection of live crisis lifelines (iCall, Vandrevala)."
            ]
        },
        {
            "title": "The AMD Slingshot Advantage",
            "content": [
                "Why AMD architecture is critical to our mission:",
                "  • On-Device ML (Ryzen AI): Shifting FER and STT models entirely to the NPU edge for absolute privacy and zero latency.",
                "  • Ollama Fallback Loop: If network drops, inference seamlessly shifts to local AMD hardware without interrupting the session.",
                "  • Datacenter Scale (EPYC): Managing thousands of concurrent duplex WebSockets with unmatched performance-per-watt."
            ]
        },
        {
            "title": "What's Next for EmpathEase?",
            "content": [
                "The road to ubiquitous, invisible support:",
                "  • Wearable Biometrics: Fusing HRV (Heart Rate Variability) into our emotional state tensor.",
                "  • Next-Gen AMD Optimization: Deep ONNX Runtime integration for hyper-efficient local NPU inference.",
                "  • Clinician Copilot: Generating multimodal summaries for human therapists to accelerate real treatment."
            ]
        },
        {
            "title": "The Edge of Empathy",
            "content": [
                "Technology shouldn't replace human connection. It should bridge the gap until human connection is possible.",
                "",
                "EmpathEase v1.1",
                "Built for AMD Slingshot 2026",
                "contact@empathease.ai"
            ]
        }
    ]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1]) 
        apply_dark_theme(slide)
        
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        format_title(title_shape, slide_data["title"])
        format_body(body_shape, slide_data["content"])
        
        # Thicker, bolder accent line
        line = shapes.add_shape(
            9, 
            title_shape.left, 
            title_shape.top + title_shape.height - Pt(5), 
            title_shape.width, 
            0
        )
        line.line.color.rgb = RGBColor(78, 205, 196) # Teal
        line.line.width = Pt(4)

    output_path = "EmpathEase_Pitch_Deck_Slingshot.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
